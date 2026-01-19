import os
import fnmatch
from pathlib import Path
from pydantic import BaseModel, Field
from . import tool, validate_args

class FilePathArgs(BaseModel):
    filepath: str = Field(..., description="The path to the file.")

class DirectoryArgs(BaseModel):
    directory: str = Field(".", description="The path to the directory.")

class WriteFileArgs(BaseModel):
    filepath: str = Field(..., description="The destination path.")
    content: str = Field(..., description="The text content to write.")

class SearchFilesArgs(BaseModel):
    directory: str = Field(..., description="Directory to search in.")
    pattern: str = Field(..., description="Search pattern (supports * and ? wildcards).")
    recursive: bool = Field(True, description="Whether to search recursively.")

class FindInFilesArgs(BaseModel):
    directory: str = Field(..., description="Directory to search in.")
    search_text: str = Field(..., description="Text to search for.")
    file_pattern: str = Field("*.py", description="File pattern to search within.")

@tool
@validate_args(DirectoryArgs)
def list_files(directory: str = ".") -> str:
    """
    Lists all files and directories in the specified path.

    Args:
        directory: The path to list. Defaults to current directory.

    Returns:
        str: A newline-separated list of items or an error message.
    """
    try:
        path = Path(directory)
        if not path.exists():
            return f"Error: Directory '{directory}' does not exist."
        if not path.is_dir():
            return f"Error: '{directory}' is not a directory."

        items = os.listdir(directory)
        return "\n".join(items) if items else "(Empty Directory)"
    except PermissionError:
        return f"Error: Permission denied for directory '{directory}'."

@tool
@validate_args(FilePathArgs)
def read_file(filepath: str) -> str:
    """
    Reads the full content of a file.

    Args:
        filepath: The path to the file to read.

    Returns:
        str: File content or an error message.
    """
    try:
        path = Path(filepath)
        if not path.exists():
            return f"Error: File '{filepath}' does not exist."
        if not path.is_file():
            return f"Error: '{filepath}' is not a file."

        # Performance: Check file size before reading (10MB limit for safety)
        file_size = path.stat().st_size
        if file_size > 10 * 1024 * 1024:
            return f"Error: File '{filepath}' is too large ({file_size / 1024 / 1024:.2f} MB). Use search or streaming tools."

        with open(path, encoding="utf-8", errors="replace") as f:
            return f.read()
    except PermissionError:
        return f"Error: Permission denied for file '{filepath}'."

@tool
@validate_args(WriteFileArgs)
def write_file(filepath: str, content: str) -> str:
    """
    Writes content to a file. Overwrites the file if it exists.

    Args:
        filepath: The destination path.
        content: The text content to write.

    Returns:
        str: Success or error message.
    """
    try:
        path = Path(filepath)
        path.parent.mkdir(parents=True, exist_ok=True)

        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"Successfully wrote {len(content)} characters to '{filepath}'."
    except PermissionError:
        return f"Error: Permission denied writing to '{filepath}'."

@tool
@validate_args(SearchFilesArgs)
def search_files(directory: str, pattern: str, recursive: bool = True) -> str:
    """
    Searches for files matching a pattern.

    Args:
        directory: Directory to search in.
        pattern: Search pattern (supports * and ? wildcards).
        recursive: Whether to search recursively.

    Returns:
        str: List of found files or message.
    """
    try:
        matches = []
        if not Path(directory).exists():
            return f"Error: Directory '{directory}' not found."
        if recursive:
            for root, _, files in os.walk(directory):
                for file in files:
                    if fnmatch.fnmatch(file, pattern):
                        matches.append(os.path.join(root, file))
        else:
            for file in os.listdir(directory):
                if fnmatch.fnmatch(file, pattern):
                    matches.append(os.path.join(directory, file))
        return (
            f"Found {len(matches)} files:\n" + "\n".join(matches[:50])
            if matches
            else "No files found matching the pattern."
        )
    except Exception as e:
        return f"Error searching files: {str(e)}"

@tool
@validate_args(FindInFilesArgs)
def find_in_files(directory: str, search_text: str, file_pattern: str = "*.py") -> str:
    """
    Searches for text within files.

    Args:
        directory: Directory to search in.
        search_text: Text to search for.
        file_pattern: File pattern to search within (e.g., *.py).

    Returns:
        str: Found results or message.
    """
    try:
        results = []
        if not Path(directory).exists():
            return f"Error: Directory '{directory}' not found."
        for root, _, files in os.walk(directory):
            for file in files:
                if fnmatch.fnmatch(file, file_pattern):
                    filepath = os.path.join(root, file)
                    try:
                        with open(filepath, encoding="utf-8", errors="ignore") as f:
                            if search_text in f.read():
                                results.append(f"{filepath}: Contains '{search_text}'")
                    except:
                        continue
        return (
            f"Found in {len(results)} files:\n" + "\n".join(results[:20])
            if results
            else f"Text '{search_text}' not found in any {file_pattern} files."
        )
    except Exception as e:
        return f"Error searching in files: {str(e)}"

@tool
@validate_args(FilePathArgs)
def read_pdf(filepath: str) -> str:
    """
    Extracts text and tables from a PDF file.

    Args:
        filepath: The path to the PDF file.

    Returns:
        str: Extracted content or error message.
    """
    try:
        import pdfplumber
        output = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                output.append(f"--- Page {i+1} ---")
                text = page.extract_text()
                if text:
                    output.append(text)
                
                tables = page.extract_tables()
                for table_idx, table in enumerate(tables):
                    output.append(f"\n[Table {table_idx+1}]")
                    for row in table:
                        # Filter out None values and join with |
                        row_str = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        output.append(row_str)
        
        return "\n".join(output) if output else "No content found in PDF."
    except ImportError:
        # Fallback to fitz if pdfplumber is not available
        try:
            import fitz
            doc = fitz.open(filepath)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text if text else "No text found in PDF (Fallback)."
        except Exception as e:
            return f"Error reading PDF (Fallback): {str(e)}"
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

@tool
@validate_args(FilePathArgs)
def read_docx(filepath: str) -> str:
    """
    Extracts text from a DOCX file.

    Args:
        filepath: The path to the DOCX file.

    Returns:
        str: Extracted text or error message.
    """
    try:
        from docx import Document
        doc = Document(filepath)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return "\n".join(text) if text else "No text found in DOCX."
    except Exception as e:
        return f"Error reading DOCX: {str(e)}"

@tool
@validate_args(FilePathArgs)
def read_excel(filepath: str) -> str:
    """
    Extracts data from an Excel file (XLSX, XLS).

    Args:
        filepath: The path to the Excel file.

    Returns:
        str: Extracted data in CSV-like format or error message.
    """
    try:
        import pandas as pd
        xl = pd.ExcelFile(filepath)
        output = []
        for sheet_name in xl.sheet_names:
            output.append(f"--- Sheet: {sheet_name} ---")
            df = xl.parse(sheet_name)
            output.append(df.to_string(index=False))
        return "\n".join(output)
    except Exception as e:
        return f"Error reading Excel: {str(e)}"

@tool
@validate_args(FilePathArgs)
def read_pptx(filepath: str) -> str:
    """
    Extracts text and slide notes from a PowerPoint file.

    Args:
        filepath: The path to the PowerPoint file.

    Returns:
        str: Extracted content or error message.
    """
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        output = []
        for i, slide in enumerate(prs.slides):
            output.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    output.append(shape.text)
            
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes:
                    output.append(f"\n[Notes]: {notes}")
        
        return "\n".join(output) if output else "No content found in PPTX."
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"
